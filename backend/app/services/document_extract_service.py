"""
文档文本抽取服务。

支持：TXT/Markdown/PDF/DOCX/PPTX
文档参考：docs/学生自建课程接入现有课程主链路_详细技术实现文档.md Section 10.5
"""

import os


class DocumentExtractService:
    """从上传文件中提取文本内容"""

    # 允许的文件后缀
    ALLOWED_SUFFIXES = {".txt", ".md", ".pdf", ".docx", ".pptx"}

    # MIME 类型映射
    MIME_MAP = {
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }

    MAX_TEXT_LENGTH = 500_000  # 最大抽取文本长度（50万字符）

    def extract_text(self, file_path: str, mime_type: str, file_name: str) -> str:
        """
        根据文件类型提取文本。

        返回：提取到的文本字符串。失败时抛出异常。
        """
        suffix = os.path.splitext(file_name)[1].lower()

        if suffix in (".txt", ".md"):
            return self._extract_text_file(file_path)
        elif suffix == ".pdf":
            return self._extract_pdf(file_path)
        elif suffix == ".docx":
            return self._extract_docx(file_path)
        elif suffix == ".pptx":
            return self._extract_pptx(file_path)
        else:
            raise ValueError(f"不支持的文件类型：{suffix}")

    def _extract_text_file(self, file_path: str) -> str:
        """提取 TXT/Markdown 文件文本"""
        # 先尝试 UTF-8，失败尝试 GBK
        for encoding in ("utf-8", "gbk", "latin-1"):
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    text = f.read()
                return text[: self.MAX_TEXT_LENGTH]
            except UnicodeDecodeError:
                continue
        raise ValueError("无法解码文件，请确保文件编码为 UTF-8 或 GBK")

    def _extract_pdf(self, file_path: str) -> str:
        """提取 PDF 文本"""
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ImportError("pypdf 未安装，无法处理 PDF 文件。请运行：pip install pypdf")

        reader = PdfReader(file_path)
        pages = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                pages.append(page_text)
            # 限制总长度
            if sum(len(p) for p in pages) >= self.MAX_TEXT_LENGTH:
                break

        text = "\n\n".join(pages)
        return text[: self.MAX_TEXT_LENGTH]

    def _extract_docx(self, file_path: str) -> str:
        """提取 DOCX 文本"""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx 未安装，无法处理 DOCX 文件。请运行：pip install python-docx")

        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
            if sum(len(p) for p in paragraphs) >= self.MAX_TEXT_LENGTH:
                break

        return "\n".join(paragraphs)[: self.MAX_TEXT_LENGTH]

    def _extract_pptx(self, file_path: str) -> str:
        """提取 PPTX 文本"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx 未安装，无法处理 PPTX 文件。请运行：pip install python-pptx")

        prs = Presentation(file_path)
        slides_text = []
        for slide in prs.slides:
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_parts.append(para.text)
            if slide_parts:
                slides_text.append("\n".join(slide_parts))
            if sum(len(s) for s in slides_text) >= self.MAX_TEXT_LENGTH:
                break

        return "\n\n---\n\n".join(slides_text)[: self.MAX_TEXT_LENGTH]

    def is_allowed_file(self, file_name: str) -> bool:
        """检查文件后缀是否允许"""
        suffix = os.path.splitext(file_name)[1].lower()
        return suffix in self.ALLOWED_SUFFIXES

    def get_mime_type(self, file_name: str) -> str:
        """根据文件后缀获取 MIME 类型"""
        suffix = os.path.splitext(file_name)[1].lower()
        return self.MIME_MAP.get(suffix, "application/octet-stream")


document_extract_service = DocumentExtractService()
